import io
from typing import List, Dict, Any
from datetime import datetime
try:
    import pptx
except ImportError:
    pptx = None

from src.parsers.base import BaseParser
from src.core.logger import setup_logger

logger = setup_logger(__name__)

class PptxParser(BaseParser):
    def parse(self, raw_data: bytes, metadata: Dict[str, Any]) -> List[Dict[str, Any]]:
        """Parse PowerPoint (.pptx) bytes."""
        if pptx is None:
            logger.error("python-pptx package is not installed.")
            return []

        source_name = metadata.get("source", "Unknown PowerPoint File")
        timestamp = datetime.now().isoformat()
        pages = []
        try:
            file_like = io.BytesIO(raw_data)
            prs = pptx.Presentation(file_like)
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text_parts.append(shape.text.strip())
                        
                slide_text = "\n".join(slide_text_parts)
                if slide_text.strip():
                    pages.append({
                        "text": slide_text,
                        "metadata": {
                            "source": source_name,
                            "timestamp": timestamp,
                            "slide_number": slide_idx + 1
                        }
                    })
            logger.info(f"Successfully extracted {len(pages)} slides from {source_name}")
        except Exception as e:
            logger.error(f"Error parsing PowerPoint {source_name}: {e}")
        return pages
